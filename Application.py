import sys, os
from PyQt5.QtWidgets import (QApplication, QMainWindow, QFileDialog, QLineEdit, QPushButton, QVBoxLayout, QHBoxLayout,
 QWidget, QPlainTextEdit, QMessageBox, QLabel, QProgressBar, QCheckBox)
from PyQt5.QtCore import Qt, pyqtSignal

from text_parser import TextReader 
from pptx import Presentation
from pptx.util import Inches

class MyWindow(QMainWindow):

    def __init__(self):
        super().__init__()
        self.delimiter = ""
        self.amounthOfSlides = None
        self.text_edit = QPlainTextEdit(self)
        self.setCentralWidget(self.text_edit)

        # перенаправление sys.stderr в self.text_edit
        sys.stderr = self.text_edit

        self.initUI()
        
    def initUI(self):
        # Создаем главный виджет и вертикальный контейнер для кнопок
        main_widget = QWidget(self)

        hbox_main = QHBoxLayout()

        vbox1 = QVBoxLayout()
        vbox2 = QVBoxLayout()

        hbox1 = QHBoxLayout()
        hbox2 = QHBoxLayout()

        hbox_main.addLayout(vbox1)
        hbox_main.addLayout(vbox2)

        vbox1.addLayout(hbox1)
        vbox1.addLayout(hbox2)

        # Создаем первую кнопку "Browse" и поле ввода для первого файла
        self.button1 = QPushButton('Read', self)
        self.button1.clicked.connect(self.on_button1_clicked)
        self.lineedit1 = QLineEdit(self)
        self.lineedit1.setReadOnly(True)

        # Создаем вторую кнопку "Browse" и поле ввода для второго файла
        self.button2 = QPushButton('Write', self)
        self.button2.clicked.connect(self.on_button2_clicked)
        self.lineedit2 = QLineEdit(self)
        self.lineedit2.setReadOnly(True)

        # Создаем "Convert"
        self.button3 = QPushButton('Convert', self)
        self.button3.clicked.connect(self.on_button3_clicked)

        # Добавляем кнопки и поля ввода в контейнер
        hbox1.addWidget(self.button1)
        hbox1.addWidget(self.lineedit1)
        hbox2.addWidget(self.button2)
        hbox2.addWidget(self.lineedit2)
        vbox1.addWidget(self.button3)

        # Добавляем ввод разделительного символа
        self.label3 = QLabel(self)
        self.label3.setAutoFillBackground(True)
        self.label3.setText("Введите разделительный символ для текста")
        self.checkbox1 = QCheckBox("Использовать ENTER")
        self.checkbox1.setChecked(True)
        self.checkbox1.clicked.connect(self.on_checkbox1_clicked)
        self.lineedit3 = QLineEdit(self)
        self.lineedit3.setStyleSheet("background-color: #f5f5f5")
        self.lineedit3.setReadOnly(True)
        vbox2.addWidget(self.label3)
        vbox2.addWidget(self.checkbox1)
        vbox2.addWidget(self.lineedit3)

        self.button4 = QPushButton('Применить символ', self)
        self.button4.clicked.connect(self.on_button4_clicked)

        vbox2.addWidget(self.button4)

        # Создаем progress bar, отображающую процесс конвертации
        self.progress_bar = QProgressBar()
        self.progress_bar.setGeometry(30, 40, 200, 25)
        self.progress_bar.setValue(0)
        vbox2.addWidget(self.progress_bar)

        # Устанавливаем контейнер как главный виджет окна
        main_widget.setLayout(hbox_main)
        self.setCentralWidget(main_widget)

        # Устанавливаем заголовок окна и размеры
        self.setWindowTitle('Text to Power Point converter')
        self.setGeometry(300, 300, 1200, 200)
        
    def on_button1_clicked(self):
        # Открываем диалог выбора файла для первого файла
        self.textFilePath, _ = QFileDialog.getOpenFileName(self, 'Open file', '')
        if self.textFilePath:
            file_extension = os.path.splitext(self.textFilePath)[1]

            if file_extension != ".txt" and file_extension != ".doc" and file_extension != ".docx":
                msg = QMessageBox()
                msg.setWindowTitle("Information")
                msg.setText("Файл с данным расширением несовместим для конвертации! Выберите другой файл")
                msg.setIcon(QMessageBox.Information)

                ok_button = msg.addButton(QMessageBox.Ok)
                msg.exec_()
            else:
                # Обновляем текст в поле ввода для первого файла
                self.lineedit1.setText(self.textFilePath)


    def on_button2_clicked(self):
        # Открываем диалог выбора файла для второго файла
        self.presFilePath, _ = QFileDialog.getSaveFileName(self, "Сохранить презентацию", "", "PowerPoint Presentation (*.pptx)")
        if self.presFilePath:
            # Обновляем текст в поле ввода для второго файла
            self.lineedit2.setText(self.presFilePath)

    def on_button3_clicked(self):
        self.text = None
        
        if self.delimiter == "":
            msg = QMessageBox()
            msg.setWindowTitle("Information")
            msg.setText("Введите разделительный символ")
            msg.setIcon(QMessageBox.Information)

            ok_button = msg.addButton(QMessageBox.Ok)
            msg.exec_()
            return

        # Проверяем наличие выбранных путей к файлам
        if self.lineedit1.text() != "" and self.lineedit2.text() != "":
            reader = TextReader(self.lineedit1.text(), self.delimiter)
            self.text = reader.read()
            self.amounthOfSlides = reader.amounthOfSlides()
            self.createPresentation()
        else:
            # Создаем окно сообщения
            msg = QMessageBox()
            msg.setWindowTitle("Information")
            msg.setText("Выберите файлы")
            msg.setIcon(QMessageBox.Information)

            # Добавляем кнопку "OK" и показываем окно
            ok_button = msg.addButton(QMessageBox.Ok)
            msg.exec_()

    def on_button4_clicked(self):
        if self.checkbox1.isChecked() == True:
            self.delimiter = "\n"
        else:
            self.delimiter = self.lineedit3.text()

    def on_checkbox1_clicked(self):
        if self.checkbox1.isChecked() == True:
            self.lineedit3.setReadOnly(True)
            self.lineedit3.setStyleSheet("background-color: #f5f5f5")
        else:
            self.lineedit3.setReadOnly(False)
            self.lineedit3.setStyleSheet("background-color: #ffffff")

    def createPresentation(self):
        presentation = Presentation()

        for slide_index in range(0, self.amounthOfSlides):
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            subtitle = slide.placeholders[1]
            subtitle.text = self.text[slide_index]
            self.progress_bar.setValue(int(((slide_index + 1) / self.amounthOfSlides) * 100))

        presentation.save(self.lineedit2.text())

        # Создаем окно сообщения
        msg = QMessageBox()
        msg.setWindowTitle("Information")
        msg.setText("Конвертация завершена")
        msg.setIcon(QMessageBox.Information)

        # Добавляем кнопку "OK" и показываем окно
        ok_button = msg.addButton(QMessageBox.Ok)
        msg.exec_()

        self.progress_bar.setValue(0)

def main():
    app = QApplication(sys.argv)
    window = MyWindow()
    window.show()
    sys.exit(app.exec_())

if __name__ == '__main__':
    main()